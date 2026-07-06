"""Offline, one-time parser: raw challenge docs -> cached backend/data/* fixtures.

Run manually (`python backend/scripts/parse_docs.py`) whenever the source
documents change. The running service never re-parses Office files at
request time (PLAN.md §3a) — it only loads the outputs of this script.
"""
import csv
import json
from pathlib import Path

import openpyxl
from docx import Document
from pptx import Presentation

ROOT = Path(__file__).resolve().parent.parent.parent
DATA_DIR = ROOT / "backend" / "data"
DATA_DIR.mkdir(exist_ok=True)


def parse_service_description() -> None:
    prs = Presentation(ROOT / "Service_description.pptx")
    lines: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        lines.append(f"--- slide {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(r.text for r in para.runs)
                    if text.strip():
                        lines.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    lines.append(" | ".join(c.text for c in row.cells))
    (DATA_DIR / "service_description.txt").write_text("\n".join(lines), encoding="utf-8")


def parse_pricing_tables() -> None:
    wb = openpyxl.load_workbook(ROOT / "PL_Industry_Challenge.xlsx", data_only=True)
    weight_bands = [
        "0-0.25kg", "0.25-0.5kg", "0.5-0.75kg", "0.75-1kg", "1-1.5kg", "1.5-2kg",
        "2-2.5kg", "2.5-3kg", "3-4kg", "4-5kg", "5-6kg", "6-7kg", "7-9kg", "9-12kg",
        "12-15kg", "15-18kg", "18-21kg", "21-24kg", "24-27kg", "27-30kg",
    ]
    volume_tiers = [
        "0-200", "201-300", "301-400", "401-500", "501-700", "701-900",
        "901-1200", "1201-1500", "1501-2000", "2001-3000", "3001-4000", "4000+",
    ]

    def read_cost_sheet(sheet_name: str) -> dict:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        data_rows = [r for r in rows if r[0] in volume_tiers]
        table = {}
        for row in data_rows:
            tier = row[0]
            table[tier] = {band: row[i + 1] for i, band in enumerate(weight_bands)}
        return table

    result = {
        "weight_bands": weight_bands,
        "volume_tiers": volume_tiers,
        "first_mile_cost": read_cost_sheet("First Mile Cost"),
        "middle_mile_cost": read_cost_sheet("Middle Mile Cost"),
        "home_delivery_cost": read_cost_sheet("Home delivery Cost"),
        "fixed_overhead_usd_per_parcel": 0.17,
        "usd_to_eur_fx_rate": 1.16,
        "region_multiplier": {"peninsula": 1.0, "balearic_islands": 1.35},
        "premium_addons_eur": {"otp": 0.35, "sod": 0.10},
        "guardrails": {
            "minimum_contribution_margin_pct": 13.0,
            "target_contribution_margin_pct": 21.0,
            "vp_approval_required_below_pct": 13.0,
            "automatic_no_go_below_pct": 9.0,
        },
    }
    (DATA_DIR / "pricing_tables.json").write_text(json.dumps(result, indent=2), encoding="utf-8")


def parse_historical_opportunities() -> None:
    wb = openpyxl.load_workbook(ROOT / "Historical_Opportunities.xlsx", data_only=True)
    ws = wb["Historical Opportunities"]
    rows = list(ws.iter_rows(values_only=True))
    header = rows[0]
    field_names = [
        "opportunity_id", "company_name", "industry", "year", "source", "incumbent_carrier",
        "daily_volume_total", "geo_fit_pct", "daily_volume_serviceable", "avg_weight_kg",
        "oversized_pct", "requires_intl", "intl_volume_share", "requires_pudo", "requires_b2b",
        "weekend_need", "annual_revenue_potential_eur", "main_pain_point", "pain_severity",
        "price_vs_incumbent_pct", "competitive_intensity", "sales_cycle_touches",
        "decision_time_days", "contract_length_months", "outcome", "lost_reason", "final_margin_pct",
    ]
    assert len(field_names) == len(header), (len(field_names), len(header))

    with (DATA_DIR / "historical_opportunities.csv").open("w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerow(field_names)
        for row in rows[1:]:
            writer.writerow(["" if v is None else v for v in row])


def parse_docx_to_text(filename: str, out_name: str) -> None:
    doc = Document(ROOT / filename)
    lines: list[str] = []
    for p in doc.paragraphs:
        if p.text.strip():
            lines.append(p.text)
    for t in doc.tables:
        for row in t.rows:
            lines.append(" | ".join(c.text for c in row.cells))
    (DATA_DIR / out_name).write_text("\n".join(lines), encoding="utf-8")


def main() -> None:
    parse_service_description()
    parse_pricing_tables()
    parse_historical_opportunities()
    parse_docx_to_text("Opportunity1_Tecnomania.docx", "tecnomania.txt")
    parse_docx_to_text("Opportunity2_PinkPapaya.docx", "pink_papaya.txt")
    print("Wrote fixtures to", DATA_DIR)


if __name__ == "__main__":
    main()
